"""
POST /api/extract-slides
  Accepts a .pptx file.
  Returns per-slide content: title, body lines, all text, a base64 thumbnail.
  Used by the frontend script editor so users can write/edit narration per slide.

POST /api/generate-script
  Accepts slide content + AI provider settings.
  Streams generated Vietnamese narration via SSE (one event per slide).
"""
import asyncio
import base64
import io
import json
import logging
from pathlib import Path
from typing import AsyncGenerator

from fastapi import APIRouter, File, Form, HTTPException, UploadFile
from fastapi.responses import StreamingResponse
from pydantic import BaseModel

logger = logging.getLogger(__name__)
router = APIRouter()


# ─────────────────────────────────────────────────────────────────────────────
# Models
# ─────────────────────────────────────────────────────────────────────────────

class SlideContent(BaseModel):
    index: int
    title: str
    body: list[str]          # bullet lines
    allText: str             # raw concatenated text for AI context
    thumbnail: str           # base64 PNG (small preview)
    hasPicture: bool


class ExtractResponse(BaseModel):
    slides: list[SlideContent]
    totalSlides: int


# ─────────────────────────────────────────────────────────────────────────────
# Extract endpoint
# ─────────────────────────────────────────────────────────────────────────────

@router.post("/extract-slides", response_model=ExtractResponse)
async def extract_slides(pptx: UploadFile = File(...)):
    """Extract text content + thumbnail from each PPTX slide."""
    if not pptx.filename or not pptx.filename.lower().endswith(".pptx"):
        raise HTTPException(400, "File phải có đuôi .pptx")

    data = await pptx.read()
    if not data:
        raise HTTPException(400, "File rỗng")

    loop = asyncio.get_event_loop()
    try:
        slides = await loop.run_in_executor(None, _extract_sync, data)
    except Exception as e:
        logger.exception("extract-slides failed")
        raise HTTPException(500, f"Không thể đọc PPTX: {e}")

    return ExtractResponse(slides=slides, totalSlides=len(slides))


def _extract_sync(pptx_bytes: bytes) -> list[SlideContent]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from services.pptx_service import (
        _convert_pillow_hifi, _find_libreoffice, CANVAS_W, CANVAS_H,
    )
    import tempfile, os

    prs = Presentation(io.BytesIO(pptx_bytes))
    results: list[SlideContent] = []

    # Render thumbnails via the same hifi renderer (small size)
    with tempfile.TemporaryDirectory() as tmpdir:
        pptx_path = Path(tmpdir) / "input.pptx"
        pptx_path.write_bytes(pptx_bytes)
        out_dir = Path(tmpdir) / "thumbs"
        out_dir.mkdir()

        # Always use Pillow for thumbnails (fast, no LibreOffice for previews)
        thumb_paths = _convert_pillow_hifi(pptx_path, out_dir, lambda m: None)

        for i, (slide, thumb_path) in enumerate(
            zip(prs.slides, thumb_paths), start=1
        ):
            title, body, all_text, has_pic = _parse_slide_text(slide)

            # Encode thumbnail as small base64 PNG (320×180)
            from PIL import Image
            with Image.open(thumb_path) as img:
                img_small = img.resize((320, 180), Image.LANCZOS)
                buf = io.BytesIO()
                img_small.save(buf, "PNG", optimize=True)
                thumb_b64 = base64.b64encode(buf.getvalue()).decode()

            results.append(SlideContent(
                index=i,
                title=title,
                body=body,
                allText=all_text,
                thumbnail=thumb_b64,
                hasPicture=has_pic,
            ))

    return results


def _parse_slide_text(slide) -> tuple[str, list[str], str, bool]:
    """Return (title, body_lines, all_text, has_picture) for a slide."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    title = ""
    body: list[str] = []
    has_pic = False

    for shape in slide.shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            has_pic = True

        if not getattr(shape, "has_text_frame", False):
            continue

        tf = shape.text_frame
        shape_text = tf.text.strip()
        if not shape_text:
            continue

        # Safely check if this shape is a title placeholder (idx==0)
        # python-pptx raises ValueError on .placeholder_format for non-placeholders
        is_title = False
        try:
            ph_fmt = shape.placeholder_format  # raises ValueError if not a placeholder
            if ph_fmt is not None and ph_fmt.idx == 0:
                is_title = True
        except (ValueError, AttributeError):
            pass

        if is_title and not title:
            title = shape_text
        else:
            for para in tf.paragraphs:
                line = para.text.strip()
                if line:
                    body.append(line)

    if not title and body:
        title = body.pop(0)

    all_text = title
    if body:
        all_text += "\n" + "\n".join(body)

    return title, body, all_text.strip(), has_pic


# ─────────────────────────────────────────────────────────────────────────────
# AI script generation endpoint
# ─────────────────────────────────────────────────────────────────────────────

@router.post("/generate-script")
async def generate_script(
    slides_json: str = Form(...),   # JSON array of {index, title, body, allText}
    provider: str = Form("gemini"), # "gemini" | "openai"
    api_key: str = Form(...),
    model: str = Form(""),          # optional model override
    style: str = Form("natural"),   # "natural" | "formal" | "friendly"
):
    """
    Generate Vietnamese narration scripts for all slides via AI.
    Streams SSE events: one per slide + a final 'done' event.

    Event format:
      data: {"index": 1, "script": "...", "done": false}
      data: {"index": -1, "script": "", "done": true}
    """
    try:
        slides = json.loads(slides_json)
    except Exception:
        raise HTTPException(400, "slides_json không hợp lệ")

    if not api_key.strip():
        raise HTTPException(400, "API key không được để trống")

    async def stream() -> AsyncGenerator[str, None]:
        for slide in slides:
            idx = slide.get("index", 0)
            all_text = slide.get("allText", "")
            title = slide.get("title", "")

            try:
                script = await _generate_one_slide(
                    index=idx,
                    title=title,
                    all_text=all_text,
                    provider=provider,
                    api_key=api_key,
                    model=model,
                    style=style,
                )
            except Exception as e:
                logger.warning(f"[AI] Slide {idx} generation failed: {e}")
                script = all_text  # fall back to raw slide text

            event = json.dumps({"index": idx, "script": script, "done": False})
            yield f"data: {event}\n\n"
            # Small pause to avoid hammering the API
            await asyncio.sleep(0.3)

        yield f"data: {json.dumps({'index': -1, 'script': '', 'done': True})}\n\n"

    return StreamingResponse(
        stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


async def _generate_one_slide(
    index: int,
    title: str,
    all_text: str,
    provider: str,
    api_key: str,
    model: str,
    style: str,
) -> str:
    """Call the chosen AI provider and return generated Vietnamese narration."""
    style_guide = {
        "natural": "tự nhiên, gần gũi, như đang nói chuyện với khán giả",
        "formal": "trang trọng, chuyên nghiệp, phù hợp hội thảo",
        "friendly": "thân thiện, vui vẻ, dễ hiểu cho người mới",
    }.get(style, "tự nhiên")

    prompt = (
        f"Bạn là chuyên gia viết kịch bản thuyết trình tiếng Việt.\n"
        f"Viết phần lời dẫn cho slide số {index} với nội dung sau:\n\n"
        f"Tiêu đề: {title}\n"
        f"Nội dung slide:\n{all_text}\n\n"
        f"Yêu cầu:\n"
        f"- Giọng văn: {style_guide}\n"
        f"- Độ dài: 2-4 câu, phù hợp để đọc trong 15-30 giây\n"
        f"- Chỉ trả về phần lời dẫn, không thêm tiêu đề hay giải thích\n"
        f"- Viết bằng tiếng Việt"
    )

    loop = asyncio.get_event_loop()

    if provider == "openai":
        return await loop.run_in_executor(
            None, _call_openai, api_key, model or "gpt-4o-mini", prompt
        )
    elif provider == "gemini":
        return await loop.run_in_executor(
            None, _call_gemini, api_key, model or "gemini-1.5-flash", prompt
        )
    else:
        raise ValueError(f"Unknown provider: {provider}")


def _call_openai(api_key: str, model: str, prompt: str) -> str:
    """Blocking OpenAI call — runs in executor."""
    try:
        from openai import OpenAI
    except ImportError:
        raise RuntimeError(
            "Thư viện 'openai' chưa được cài. Chạy: pip install openai"
        )
    client = OpenAI(api_key=api_key)
    resp = client.chat.completions.create(
        model=model,
        messages=[{"role": "user", "content": prompt}],
        max_tokens=300,
        temperature=0.7,
    )
    return resp.choices[0].message.content.strip()


def _call_gemini(api_key: str, model: str, prompt: str) -> str:
    """Blocking Gemini call — runs in executor."""
    try:
        import google.generativeai as genai
    except ImportError:
        raise RuntimeError(
            "Thư viện 'google-generativeai' chưa được cài. "
            "Chạy: pip install google-generativeai"
        )
    genai.configure(api_key=api_key)
    m = genai.GenerativeModel(model)
    resp = m.generate_content(prompt)
    return resp.text.strip()
