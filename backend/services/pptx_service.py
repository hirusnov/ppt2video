"""
PPTX → PNG conversion service.

Priority chain:
1. LibreOffice headless  — 100% faithful (production/Docker)
2. High-fidelity Pillow  — extracts real backgrounds, images, shapes & text
                           from python-pptx (dev / no LibreOffice)
"""
import asyncio
import io
import os
import logging
import re
import shutil
from pathlib import Path
from typing import Callable, Awaitable

logger = logging.getLogger(__name__)

# Slide canvas size for Pillow renderer (16:9)
CANVAS_W = 1280
CANVAS_H = 720


# ─────────────────────────────────────────────────────────────────────────────
# Public entry point
# ─────────────────────────────────────────────────────────────────────────────

async def convert_to_images(
    pptx_path: Path,
    output_dir: Path,
    log_cb: Callable[[str], Awaitable[None]] | None = None,
) -> list[Path]:
    output_dir.mkdir(parents=True, exist_ok=True)

    async def _log(msg: str) -> None:
        logger.info(msg)
        if log_cb:
            await log_cb(msg)

    lo = _find_libreoffice()
    if lo:
        return await _convert_libreoffice(pptx_path, output_dir, lo, _log)

    await _log("[PPTX] LibreOffice not found — using high-fidelity Pillow renderer")
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(
        None, _convert_pillow_hifi, pptx_path, output_dir, _log_sync
    )


def _log_sync(msg: str) -> None:
    logger.info(msg)


# ─────────────────────────────────────────────────────────────────────────────
# LibreOffice path
# ─────────────────────────────────────────────────────────────────────────────

def _find_libreoffice() -> str | None:
    for c in [
        # Windows
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        # Linux / Docker
        "libreoffice", "soffice",
        "/usr/bin/libreoffice", "/usr/bin/soffice",
        "/usr/lib/libreoffice/program/soffice",
        "/opt/libreoffice/program/soffice",
    ]:
        if shutil.which(c) or (os.path.isfile(c) if os.path.isabs(c) else False):
            return c
    return None


async def _convert_libreoffice(
    pptx_path: Path,
    output_dir: Path,
    lo_bin: str,
    log_cb: Callable[[str], Awaitable[None]],
) -> list[Path]:
    await log_cb(f"[PPTX] Converting with LibreOffice: {pptx_path.name}")
    cmd = [
        lo_bin,
        "--headless",
        "--norestore",
        "--nofirststartwizard",
        "--nologo",
        "--nolockcheck",
        "--convert-to", "png",
        "--outdir", str(output_dir),
        str(pptx_path),
    ]

    # Use blocking subprocess.run in thread pool — avoids asyncio ProactorEventLoop
    # requirement on Windows (which breaks with uvicorn --reload).
    loop = asyncio.get_event_loop()
    returncode, stderr_text = await loop.run_in_executor(
        None, _run_libreoffice_blocking, cmd
    )

    if returncode != 0:
        raise RuntimeError(
            f"LibreOffice failed: {stderr_text[-600:]}"
        )

    raw = sorted(
        output_dir.glob("*.png"),
        key=lambda p: _lo_sort_key(p.stem),
    )
    if not raw:
        raise RuntimeError("LibreOffice produced no PNG files")

    renamed = []
    for i, src in enumerate(raw, start=1):
        dst = output_dir / f"slide_{i:03d}.png"
        src.rename(dst)
        renamed.append(dst)
        await log_cb(f"[PPTX] Converted {i}/{len(raw)} slides to PNG")

    await log_cb(f"[PPTX] Done — {len(renamed)} slides exported")
    return renamed


def _run_libreoffice_blocking(cmd: list[str]) -> tuple[int, str]:
    """Run LibreOffice as a blocking subprocess — safe on all platforms."""
    import subprocess
    try:
        result = subprocess.run(
            cmd,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=120,
        )
        return result.returncode, result.stderr.decode(errors="replace")
    except subprocess.TimeoutExpired:
        raise RuntimeError("LibreOffice timed out after 120s")


def _lo_sort_key(stem: str) -> tuple[str, int]:
    m = re.match(r"^(.*?)(\d+)$", stem)
    return (m.group(1), int(m.group(2))) if m else (stem, 0)


# ─────────────────────────────────────────────────────────────────────────────
# High-fidelity Pillow renderer
# ─────────────────────────────────────────────────────────────────────────────

def _convert_pillow_hifi(
    pptx_path: Path,
    output_dir: Path,
    log_cb: Callable[[str], None],
) -> list[Path]:
    """
    Render each PPTX slide to a PNG by:
    1. Filling the background (solid colour or background image from theme/layout)
    2. Compositing all picture shapes (embedded images) at correct positions
    3. Drawing all text frames with correct position, size, font, colour, bold/italic
    """
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.dml.color import RGBColor
    from PIL import Image, ImageDraw, ImageFont

    prs = Presentation(str(pptx_path))
    slide_w_emu = prs.slide_width   # EMU
    slide_h_emu = prs.slide_height

    def emu_to_px(emu: int, axis: str) -> int:
        if axis == "x":
            return round(emu / slide_w_emu * CANVAS_W)
        return round(emu / slide_h_emu * CANVAS_H)

    paths: list[Path] = []
    total = len(prs.slides)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        img = Image.new("RGB", (CANVAS_W, CANVAS_H), (255, 255, 255))

        # ── 1. Background fill ──────────────────────────────────────────────
        bg_color = _extract_slide_bg_color(slide, prs)
        if bg_color:
            img = Image.new("RGB", (CANVAS_W, CANVAS_H), bg_color)

        # Try background image from slide background
        bg_img = _extract_slide_bg_image(slide)
        if bg_img:
            bg_img = bg_img.resize((CANVAS_W, CANVAS_H), Image.LANCZOS)
            img.paste(bg_img, (0, 0))

        draw = ImageDraw.Draw(img)

        # ── 2. Shapes (pictures first, then text on top) ────────────────────
        shapes_sorted = sorted(
            slide.shapes,
            key=lambda s: (
                0 if getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE else 1
            ),
        )

        for shape in shapes_sorted:
            left   = emu_to_px(shape.left   or 0, "x")
            top    = emu_to_px(shape.top    or 0, "y")
            width  = emu_to_px(shape.width  or 0, "x")
            height = emu_to_px(shape.height or 0, "y")

            # ── Picture shapes ──────────────────────────────────────────────
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                try:
                    pic_bytes = shape.image.blob
                    pic_img = Image.open(io.BytesIO(pic_bytes)).convert("RGBA")
                    if width > 0 and height > 0:
                        pic_img = pic_img.resize((width, height), Image.LANCZOS)
                    # Composite with alpha
                    tmp = Image.new("RGBA", img.size, (0, 0, 0, 0))
                    tmp.paste(pic_img, (left, top))
                    img = Image.alpha_composite(img.convert("RGBA"), tmp).convert("RGB")
                    draw = ImageDraw.Draw(img)
                except Exception as e:
                    logger.debug(f"[PPTX] Could not render picture shape: {e}")
                continue

            # ── Filled rectangle / autoshape background ─────────────────────
            fill_color = _get_shape_fill_color(shape)
            if fill_color and width > 0 and height > 0:
                draw.rectangle([left, top, left + width, top + height], fill=fill_color)

            # ── Text frames ─────────────────────────────────────────────────
            if not getattr(shape, "has_text_frame", False):
                continue
            if not shape.text_frame.text.strip():
                continue

            tf = shape.text_frame
            # Word-wrap within shape bounds
            cursor_y = top + emu_to_px(tf.margin_top or 0, "y")
            max_x = left + max(width - emu_to_px(tf.margin_left or 0, "x")
                                     - emu_to_px(tf.margin_right or 0, "x"), 40)
            text_left = left + emu_to_px(tf.margin_left or 0, "x")

            for para in tf.paragraphs:
                para_text = para.text
                if not para_text:
                    cursor_y += 8
                    continue

                # Font size: use run-level, para-level, or default
                font_size = _get_para_font_size(para, default=18)
                bold = _get_para_bold(para)
                color = _get_para_color(para) or (0, 0, 0)
                align = _get_para_align(para)

                font = _load_font(font_size, bold)
                lines = _wrap_text(draw, para_text, font, max_x - text_left)

                for line in lines:
                    if cursor_y > top + height:
                        break
                    x = text_left
                    if align == "center":
                        try:
                            bbox = draw.textbbox((0, 0), line, font=font)
                            x = text_left + (max_x - text_left - (bbox[2] - bbox[0])) // 2
                        except Exception:
                            pass
                    elif align == "right":
                        try:
                            bbox = draw.textbbox((0, 0), line, font=font)
                            x = max_x - (bbox[2] - bbox[0])
                        except Exception:
                            pass

                    # Shadow for readability
                    draw.text((x + 1, cursor_y + 1), line, font=font, fill=(0, 0, 0, 80))
                    draw.text((x, cursor_y), line, font=font, fill=color)
                    cursor_y += font_size + 4

                cursor_y += 6  # paragraph spacing

        out_path = output_dir / f"slide_{slide_idx:03d}.png"
        img.save(str(out_path), "PNG", optimize=True)
        paths.append(out_path)
        log_cb(f"[PPTX] Rendered slide {slide_idx}/{total}")

    log_cb(f"[PPTX] Done — {len(paths)} slides rendered")
    return paths


# ─────────────────────────────────────────────────────────────────────────────
# Helper: background extraction
# ─────────────────────────────────────────────────────────────────────────────

def _extract_slide_bg_color(slide, prs) -> tuple[int, int, int] | None:
    """Try to get a solid background colour from the slide's background fill."""
    try:
        bg = slide.background
        fill = bg.fill
        fill.background()  # just to poke it; ignore errors
        if fill.type is None:
            return None
        from pptx.enum.dml import MSO_THEME_COLOR
        from pptx.dml.color import RGBColor
        if str(fill.type) == "SOLID (1)":
            fg = fill.fore_color
            rgb = fg.rgb
            return (rgb.r, rgb.g, rgb.b)
    except Exception:
        pass

    # Fallback: try slide layout / slide master
    for source in [slide.slide_layout, slide.slide_layout.slide_master]:
        try:
            fill = source.background.fill
            if str(fill.type) == "SOLID (1)":
                rgb = fill.fore_color.rgb
                return (rgb.r, rgb.g, rgb.b)
        except Exception:
            continue

    return None


def _extract_slide_bg_image(slide) -> "Image.Image | None":
    """Extract a background image if the slide has one."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from PIL import Image
        bg = slide.background
        fill = bg.fill
        # Picture fill
        if hasattr(fill, "_fill") and fill._fill is not None:
            blip_fill = fill._fill.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}blipFill"
            )
            if blip_fill is not None:
                blip = blip_fill.find(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
                )
                if blip is not None:
                    rId = blip.get(
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                    )
                    if rId:
                        img_part = slide.part.related_parts.get(rId)
                        if img_part:
                            return Image.open(io.BytesIO(img_part.blob)).convert("RGB")
    except Exception:
        pass
    return None


# ─────────────────────────────────────────────────────────────────────────────
# Helper: shape fill
# ─────────────────────────────────────────────────────────────────────────────

def _get_shape_fill_color(shape) -> tuple[int, int, int] | None:
    try:
        fill = shape.fill
        if str(fill.type) == "SOLID (1)":
            rgb = fill.fore_color.rgb
            return (rgb.r, rgb.g, rgb.b)
    except Exception:
        pass
    return None


# ─────────────────────────────────────────────────────────────────────────────
# Helper: text formatting
# ─────────────────────────────────────────────────────────────────────────────

def _get_para_font_size(para, default: int = 18) -> int:
    """Return font size in pixels (approx) for a paragraph."""
    # Try run-level first
    for run in para.runs:
        try:
            if run.font.size:
                return max(8, round(run.font.size.pt * 96 / 72))
        except Exception:
            pass
    # Para-level
    try:
        if para.font.size:
            return max(8, round(para.font.size.pt * 96 / 72))
    except Exception:
        pass
    return default


def _get_para_bold(para) -> bool:
    for run in para.runs:
        try:
            if run.font.bold:
                return True
        except Exception:
            pass
    try:
        return bool(para.font.bold)
    except Exception:
        return False


def _get_para_color(para) -> tuple[int, int, int] | None:
    for run in para.runs:
        try:
            rgb = run.font.color.rgb
            return (rgb.r, rgb.g, rgb.b)
        except Exception:
            pass
    try:
        rgb = para.font.color.rgb
        return (rgb.r, rgb.g, rgb.b)
    except Exception:
        pass
    return None


def _get_para_align(para) -> str:
    try:
        from pptx.enum.text import PP_ALIGN
        a = para.alignment
        if a == PP_ALIGN.CENTER:
            return "center"
        if a == PP_ALIGN.RIGHT:
            return "right"
    except Exception:
        pass
    return "left"


# ─────────────────────────────────────────────────────────────────────────────
# Helper: font loading with cache
# ─────────────────────────────────────────────────────────────────────────────

_font_cache: dict[tuple, object] = {}


def _load_font(size: int, bold: bool = False):
    from PIL import ImageFont
    key = (size, bold)
    if key in _font_cache:
        return _font_cache[key]

    # Try common system fonts in order of preference
    font_candidates: list[str] = []
    import sys, os
    if sys.platform == "win32":
        windir = os.environ.get("WINDIR", "C:\\Windows")
        fonts_dir = os.path.join(windir, "Fonts")
        for fname in (
            ["calibrib.ttf", "calibri.ttf", "arial.ttf", "arialbd.ttf",
             "segoeui.ttf", "seguisb.ttf", "tahoma.ttf", "verdana.ttf"]
        ):
            font_candidates.append(os.path.join(fonts_dir, fname))
    else:
        for fname in [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
            "/usr/share/fonts/noto/NotoSans-Bold.ttf",
            "/usr/share/fonts/noto/NotoSans-Regular.ttf",
        ]:
            font_candidates.append(fname)

    font = None
    for path in font_candidates:
        if os.path.exists(path):
            try:
                font = ImageFont.truetype(path, size)
                break
            except Exception:
                continue

    if font is None:
        try:
            font = ImageFont.load_default(size=size)
        except Exception:
            font = ImageFont.load_default()

    _font_cache[key] = font
    return font


def _wrap_text(draw, text: str, font, max_width: int) -> list[str]:
    """Word-wrap text to fit within max_width pixels."""
    if max_width <= 0:
        return [text]
    words = text.split()
    lines: list[str] = []
    current = ""
    for word in words:
        test = f"{current} {word}".strip()
        try:
            bbox = draw.textbbox((0, 0), test, font=font)
            w = bbox[2] - bbox[0]
        except Exception:
            w = len(test) * (font.size if hasattr(font, "size") else 10)
        if w > max_width and current:
            lines.append(current)
            current = word
        else:
            current = test
    if current:
        lines.append(current)
    return lines or [text]
